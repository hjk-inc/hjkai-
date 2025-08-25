import requests
from bs4 import BeautifulSoup
import json
import numpy as np
import os
import customtkinter as ctk
from tkinter import filedialog, messagebox, Toplevel, ttk
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from datetime import datetime
import time
import random
import threading
import sys
import re
import networkx as nx
import sqlite3
import hashlib
import spacy
import joblib
import webbrowser
import logging
import plotly.graph_objects as go
from transformers import pipeline
from sentence_transformers import SentenceTransformer
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from sklearn.neighbors import NearestNeighbors
from sklearn.linear_model import LogisticRegression
from sklearn.decomposition import LatentDirichletAllocation
from sklearn.metrics import silhouette_score
from sklearn.base import clone
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import RandomForestClassifier
from sklearn.manifold import TSNE
import pandas as pd
from PIL import Image, ImageTk
import pyttsx3
import speech_recognition as sr
from langdetect import detect
import importlib.util
import inspect
import base64
import io

# --- Configure logging ---
logging.basicConfig(
    filename='hjkai.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger('HJKAI')

# --- Initialize NLP ---
try:
    nlp = spacy.load("en_core_web_lg")
except:
    import spacy.cli
    spacy.cli.download("en_core_web_lg")
    nlp = spacy.load("en_core_web_lg")

# --- Global Models ---
SENTIMENT_ANALYZER = None
SUMMARIZER = None
TRANSLATOR = None
UNIVERSAL_ENCODER = None

# Initialize models in a thread-safe way
def initialize_models():
    global SENTIMENT_ANALYZER, SUMMARIZER, TRANSLATOR, UNIVERSAL_ENCODER
    try:
        SENTIMENT_ANALYZER = pipeline("sentiment-analysis")
        SUMMARIZER = pipeline("summarization")
        TRANSLATOR = pipeline("translation_en_to_fr", model="t5-small")
        UNIVERSAL_ENCODER = SentenceTransformer('all-mpnet-base-v2')
        logger.info("Models initialized successfully")
    except Exception as e:
        logger.error(f"Model initialization failed: {str(e)}")

# Run in background thread
threading.Thread(target=initialize_models, daemon=True).start()

# --- Banned Word Base ---
banned_words = ['ads', 'sponsored', 'click here']

# --- Database Setup ---
def init_db():
    conn = sqlite3.connect('hjkai.db')
    c = conn.cursor()
    
    # Queries table
    c.execute('''CREATE TABLE IF NOT EXISTS queries
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  topic TEXT,
                  timestamp DATETIME DEFAULT CURRENT_TIMESTAMP,
                  results_file TEXT)''')
    
    # Banned words
    c.execute('''CREATE TABLE IF NOT EXISTS banned_words
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  word TEXT UNIQUE,
                  added_date DATETIME DEFAULT CURRENT_TIMESTAMP)''')
    
    # Models
    c.execute('''CREATE TABLE IF NOT EXISTS models
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  model_type TEXT,
                  model_file TEXT,
                  timestamp DATETIME DEFAULT CURRENT_TIMESTAMP)''')
    
    # Knowledge graph
    c.execute('''CREATE TABLE IF NOT EXISTS knowledge_graph
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  source TEXT,
                  target TEXT,
                  relation TEXT,
                  strength REAL,
                  UNIQUE(source, target))''')
    
    # Concepts
    c.execute('''CREATE TABLE IF NOT EXISTS concepts
                 (concept TEXT PRIMARY KEY,
                  embedding BLOB,
                  last_updated DATETIME DEFAULT CURRENT_TIMESTAMP)''')
    
    conn.commit()
    conn.close()

# Initialize databases
init_db()

# --- Reasoning System ---
class AIReasoner:
    def __init__(self):
        self.knowledge_base = {
            "filtering": "Rejects content with banned words, short texts, or low relevance terms",
            "clustering": "Groups similar content using KMeans on TF-IDF vectors",
            "question_gen": "Uses nearest neighbors on past topics to generate questions",
            "self_improve": "Learns new banned words from clean/spam samples using logistic regression",
            "sentiment": "Analyzes emotional tone using Transformers",
            "summarization": "Generates abstracts using extractive and abstractive methods",
            "knowledge_graph": "Builds semantic relationships between concepts",
            "universal_mode": "Connects concepts across all domains of knowledge",
            "auto_improve": "Continuously enhances models after each query"
        }
        
    def explain_decision(self, text):
        """Provide detailed reasoning for filtering decisions with confidence score"""
        reasons = []
        confidence = 1.0  # Start with full confidence
        
        # Check banned words
        found_banned = []
        for word in banned_words:
            if word in text.lower():
                found_banned.append(word)
                confidence *= 0.7  # Reduce confidence for each banned word
        if found_banned:
            reasons.append(f"Contains banned terms: {', '.join(found_banned)}")
        
        # Check length
        word_count = len(text.split())
        if word_count < 8:
            reasons.append(f"Too short ({word_count} words < 8 minimum)")
            confidence *= 0.9
            
        # Check relevance
        relevance_terms = ['science', 'tech', 'data', 'research', 'study', 'analysis']
        found_relevance = [term for term in relevance_terms if term in text.lower()]
        if not found_relevance:
            reasons.append(f"Missing relevance terms ({', '.join(relevance_terms)})")
            confidence *= 0.8
            
        # Calculate risk level
        risk_level = "low"
        if confidence < 0.5:
            risk_level = "high"
        elif confidence < 0.7:
            risk_level = "medium"
            
        return reasons, confidence, risk_level
    
    def explain_cluster(self, cluster_id, keywords):
        """Explain cluster characteristics"""
        return f"Cluster {cluster_id} focuses on: {', '.join(keywords)}"
    
    def system_report(self):
        """Generate full system reasoning report"""
        report = "⚡ HJKAI Reasoning System\n"
        report += "="*50 + "\n"
        report += f"Banned words: {', '.join(banned_words)}\n"
        report += f"Knowledge base size: {len(self.knowledge_base)} concepts\n"
        
        # Add database stats
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        c.execute("SELECT COUNT(*) FROM queries")
        query_count = c.fetchone()[0]
        c.execute("SELECT COUNT(*) FROM banned_words")
        banned_count = c.fetchone()[0]
        c.execute("SELECT COUNT(*) FROM models")
        model_count = c.fetchone()[0]
        c.execute("SELECT COUNT(*) FROM knowledge_graph")
        kg_count = c.fetchone()[0]
        conn.close()
        
        report += f"Historical queries: {query_count}\n"
        report += f"Learned banned words: {banned_count}\n"
        report += f"Saved models: {model_count}\n"
        report += f"Knowledge graph relations: {kg_count}\n"
        report += "Capabilities:\n"
        for cap, desc in self.knowledge_base.items():
            report += f"- {cap.upper()}: {desc}\n"
        return report

# --- Core AI System ---
class HJKAICore:
    def __init__(self):
        self.reasoner = AIReasoner()
        self.question_model = self.train_or_load_question_model()
        self.knowledge_graph = nx.Graph()
        self.load_knowledge_graph()
        self.recognizer = sr.Recognizer()
        self.engine = pyttsx3.init()
        self.engine.setProperty('rate', 150)  # Speed percent
        self.engine.setProperty('volume', 0.9)  # Volume 0-1
        
    def train_or_load_question_model(self):
        """Load saved model or train new one"""
        model_file = "question_model.joblib"
        if os.path.exists(model_file):
            try:
                return joblib.load(model_file)
            except:
                logger.warning("Failed to load question model, retraining...")
                
        return self.train_question_model()
    
    def train_question_model(self):
        logger.info("Training new question model")
        past_topics = ["solar energy", "machine learning", "climate change", "quantum computing", 
                      "neural networks", "genetic engineering", "nanotechnology", "blockchain"]
        past_questions = [
            self.generate_default_questions("solar energy"),
            self.generate_default_questions("machine learning"),
            self.generate_default_questions("climate change"),
            self.generate_default_questions("quantum computing"),
            self.generate_default_questions("neural networks"),
            self.generate_default_questions("genetic engineering"),
            self.generate_default_questions("nanotechnology"),
            self.generate_default_questions("blockchain")
        ]
        vectorizer = TfidfVectorizer()
        X = vectorizer.fit_transform(past_topics)
        model = NearestNeighbors(n_neighbors=1).fit(X)
        
        # Save model
        joblib.dump((vectorizer, model, past_topics, past_questions), "question_model.joblib")
        
        # Save to database
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        c.execute("INSERT INTO models (model_type, model_file) VALUES (?, ?)", 
                 ("question_generator", "question_model.joblib"))
        conn.commit()
        conn.close()
        
        return vectorizer, model, past_topics, past_questions
    
    def generate_default_questions(self, topic):
        return [
            f"What is {topic}?",
            f"How does {topic} work?",
            f"Why is {topic} important?",
            f"What are the future uses of {topic}?",
            f"What are the main challenges in {topic}?",
            f"Who are the key researchers in {topic}?",
            f"What recent breakthroughs have occurred in {topic}?"
        ]
    
    def ml_generate_questions(self, topic):
        vec, model_nn, topics, past_questions = self.question_model
        X_topic = vec.transform([topic])
        _, idx = model_nn.kneighbors(X_topic)
        return past_questions[idx[0][0]]
    
    def upgrade_banned_words(self, clean_samples, spam_samples):
        global banned_words
        texts = clean_samples + spam_samples
        labels = [0] * len(clean_samples) + [1] * len(spam_samples)
        vectorizer = TfidfVectorizer(stop_words='english', max_features=100)
        X = vectorizer.fit_transform(texts)
        clf = LogisticRegression().fit(X, labels)
        top_indices = np.argsort(clf.coef_[0])[-5:]
        new_banned = [vectorizer.get_feature_names_out()[i] for i in top_indices]
        banned_words.extend(new_banned)
        banned_words = list(set(banned_words))
        
        # Save to database
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        for word in new_banned:
            try:
                c.execute("INSERT OR IGNORE INTO banned_words (word) VALUES (?)", (word,))
            except sqlite3.IntegrityError:
                pass  # Already exists
        conn.commit()
        conn.close()
    
    def fetch_web_data(self, topic, search_engine="google"):
        delay = random.uniform(1.5, 4.0)
        time.sleep(delay)
        
        user_agents = [
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36",
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/14.1.1 Safari/605.1.15",
            "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/92.0.4515.107 Safari/537.36"
        ]
        
        headers = {"User-Agent": random.choice(user_agents)}
        
        try:
            if search_engine == "google":
                search_url = f"https://www.google.com/search?q={topic.replace(' ', '+')}"
            elif search_engine == "bing":
                search_url = f"https://www.bing.com/search?q={topic.replace(' ', '+')}"
            elif search_engine == "duckduckgo":
                search_url = f"https://duckduckgo.com/html/?q={topic.replace(' ', '+')}"
            else:
                search_url = f"https://www.google.com/search?q={topic.replace(' ', '+')}"
            
            res = requests.get(search_url, headers=headers, timeout=10)
            res.raise_for_status()
            
            soup = BeautifulSoup(res.text, "html.parser")
            results = []
            
            # Different parsing for different search engines
            if search_engine == "google":
                for div in soup.find_all("div", class_="tF2Cxc"):
                    snippet = div.find("div", class_="VwiC3b")
                    if snippet:
                        text = snippet.get_text()
                        if len(text) > 30:
                            results.append(text)
            elif search_engine == "bing":
                for li in soup.find_all("li", class_="b_algo"):
                    snippet = li.find("p")
                    if snippet:
                        text = snippet.get_text()
                        if len(text) > 30:
                            results.append(text)
            elif search_engine == "duckduckgo":
                for div in soup.find_all("div", class_="result__snippet"):
                    text = div.get_text()
                    if len(text) > 30:
                        results.append(text)
                        
            return results[:10] if len(results) > 5 else []
        
        except Exception as e:
            logger.error(f"Scraping error: {str(e)}")
            raise
    
    def chunk_text(self, text, max_length=500):
        """Split text into chunks for summarization"""
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + 1 <= max_length:
                current_chunk.append(word)
                current_length += len(word) + 1
            else:
                chunks.append(" ".join(current_chunk))
                current_chunk = [word]
                current_length = len(word)
        
        if current_chunk:
            chunks.append(" ".join(current_chunk))
            
        return chunks
    
    def summarize_text(self, text):
        """Summarize text handling long documents"""
        if not SUMMARIZER:
            return text[:300] + "..."  # Fallback
            
        try:
            if len(text.split()) > 500:
                chunks = self.chunk_text(text)
                summaries = []
                for chunk in chunks:
                    summary = SUMMARIZER(chunk, max_length=150, min_length=30, do_sample=False)[0]['summary_text']
                    summaries.append(summary)
                return " ".join(summaries)
            else:
                return SUMMARIZER(text, max_length=150, min_length=30, do_sample=False)[0]['summary_text']
        except Exception as e:
            logger.error(f"Summarization error: {str(e)}")
            return text[:300] + "..."  # Fallback
    
    def detect_language(self, text):
        try:
            return detect(text)
        except:
            return "en"  # Default to English
    
    def translate_text(self, text, target_lang="en"):
        if not TRANSLATOR:
            return text  # Fallback
            
        try:
            src_lang = self.detect_language(text)
            if src_lang == target_lang:
                return text
                
            return TRANSLATOR(text, src_lang=src_lang, tgt_lang=target_lang)[0]['translation_text']
        except Exception as e:
            logger.error(f"Translation error: {str(e)}")
            return text
    
    def train_on_data(self, data, use_saved=True):
        if not data:
            return {}, {}, [], {}, ""
        
        # Try to load saved model
        model_hash = hashlib.md5(" ".join(data).encode()).hexdigest()
        model_file = f"cluster_model_{model_hash}.joblib"
        
        if use_saved and os.path.exists(model_file):
            try:
                return joblib.load(model_file)
            except:
                logger.warning("Failed to load cluster model, retraining...")
        
        vectorizer = TfidfVectorizer(stop_words='english', max_features=5000)
        X = vectorizer.fit_transform(data)
        
        # Determine optimal clusters
        max_clusters = min(10, len(data))
        silhouette_scores = []
        for n in range(2, max_clusters+1):
            model = KMeans(n_clusters=n, random_state=0)
            labels = model.fit_predict(X)
            if len(set(labels)) > 1:  # Silhouette requires at least 2 clusters
                score = silhouette_score(X, labels)
                silhouette_scores.append(score)
            else:
                silhouette_scores.append(0)
                
        optimal_n = np.argmax(silhouette_scores) + 2  # +2 because range started at 2
        model = KMeans(n_clusters=optimal_n, random_state=0).fit(X)
        labels = model.labels_
        
        clustered_data = {f"Cluster {i}": [] for i in set(labels)}
        for text, label in zip(data, labels):
            clustered_data[f"Cluster {label}"].append(text)

        # Extract keywords
        keywords_per_cluster = {}
        order_centroids = model.cluster_centers_.argsort()[:, ::-1]
        terms = vectorizer.get_feature_names_out()
        for i in range(model.n_clusters):
            top_words = [terms[ind] for ind in order_centroids[i, :10]]
            keywords_per_cluster[f"Cluster {i}"] = top_words

        # Topic modeling with LDA
        lda = LatentDirichletAllocation(n_components=min(3, optimal_n), 
                                        random_state=0,
                                        learning_method='online')
        lda.fit(X)
        
        # Create topic summaries
        topic_summaries = {}
        for idx, topic in enumerate(lda.components_):
            top_indices = topic.argsort()[:-11:-1]
            top_terms = [terms[i] for i in top_indices]
            topic_summaries[f"Topic {idx}"] = ", ".join(top_terms)
        
        # Generate overall summary
        try:
            full_text = " ".join(data)
            summary = self.summarize_text(full_text)
        except:
            summary = " ".join(data)[:300] + "..."  # Fallback
        
        # Save model
        result = (clustered_data, keywords_per_cluster, labels, topic_summaries, summary)
        joblib.dump(result, model_file)
        
        # Save to database
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        c.execute("INSERT INTO models (model_type, model_file) VALUES (?, ?)", 
                 ("clustering", model_file))
        conn.commit()
        conn.close()
        
        return result
    
    def analyze_sentiment(self, data):
        sentiments = []
        for text in data:
            if SENTIMENT_ANALYZER:
                try:
                    result = SENTIMENT_ANALYZER(text)[0]
                    sentiment = {
                        "text": text,
                        "sentiment": result['label'],
                        "score": result['score']
                    }
                except:
                    # Fallback to TextBlob
                    from textblob import TextBlob
                    blob = TextBlob(text)
                    sentiment = {
                        "text": text,
                        "sentiment": "positive" if blob.sentiment.polarity > 0.1 else 
                                     "negative" if blob.sentiment.polarity < -0.1 else "neutral",
                        "score": abs(blob.sentiment.polarity)
                    }
            else:
                # Fallback to TextBlob
                from textblob import TextBlob
                blob = TextBlob(text)
                sentiment = {
                    "text": text,
                    "sentiment": "positive" if blob.sentiment.polarity > 0.1 else 
                                 "negative" if blob.sentiment.polarity < -0.1 else "neutral",
                    "score": abs(blob.sentiment.polarity)
                }
            sentiments.append(sentiment)
        
        # Create sentiment distribution
        sentiment_counts = {
            "positive": sum(1 for s in sentiments if s["sentiment"] == "POSITIVE" or s["sentiment"] == "positive"),
            "neutral": sum(1 for s in sentiments if s["sentiment"] == "NEUTRAL" or s["sentiment"] == "neutral"),
            "negative": sum(1 for s in sentiments if s["sentiment"] == "NEGATIVE" or s["sentiment"] == "negative")
        }
        
        return sentiments, sentiment_counts
    
    def build_knowledge_graph(self, data, keywords):
        G = nx.Graph()
        
        # Add keywords as nodes
        for cluster, terms in keywords.items():
            for term in terms[:5]:
                G.add_node(term, type="concept", cluster=cluster)
        
        # Add connections based on co-occurrence
        for text in data:
            # Simple implementation - real version would use NLP
            words = text.lower().split()
            for term in G.nodes:
                if term in words:
                    for other_term in G.nodes:
                        if other_term != term and other_term in words:
                            if G.has_edge(term, other_term):
                                G[term][other_term]["weight"] += 1
                            else:
                                G.add_edge(term, other_term, weight=1)
    
        return G
    
    def auto_improve(self, filtered_data):
        """Automatically improve models after each query"""
        # Update question model with new topic
        # (In real implementation, we'd add this topic to training data)
        logger.info("Auto-improving models...")
        
        # Save updated banned words
        with open("banned_words.json", "w") as f:
            json.dump(banned_words, f)
        
        # Update clustering model
        if filtered_data:
            try:
                self.train_on_data(filtered_data, use_saved=False)
            except Exception as e:
                logger.error(f"Auto-improve error: {str(e)}")
    
    def load_knowledge_graph(self):
        """Load knowledge graph from database"""
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        
        # Load concepts
        c.execute("SELECT concept, embedding FROM concepts")
        for concept, embedding_blob in c.fetchall():
            embedding = np.frombuffer(embedding_blob, dtype=np.float32)
            self.knowledge_graph.add_node(concept, embedding=embedding)
        
        # Load relations
        c.execute("SELECT source, target, relation, strength FROM knowledge_graph")
        for source, target, rel_type, strength in c.fetchall():
            self.knowledge_graph.add_edge(source, target, relation=rel_type, strength=strength)
        
        conn.close()
        logger.info(f"Loaded knowledge graph with {len(self.knowledge_graph.nodes)} nodes")
    
    def update_knowledge_graph(self, topic, related_topics):
        """Update universal knowledge graph with new information"""
        # Get or create topic node
        if topic not in self.knowledge_graph:
            if UNIVERSAL_ENCODER:
                topic_embedding = UNIVERSAL_ENCODER.encode([topic])[0]
            else:
                topic_embedding = np.random.rand(768)  # Fallback
            self.knowledge_graph.add_node(topic, embedding=topic_embedding)
            
            # Save to database
            conn = sqlite3.connect('hjkai.db')
            c = conn.cursor()
            c.execute("INSERT OR REPLACE INTO concepts (concept, embedding) VALUES (?, ?)", 
                     (topic, topic_embedding.tobytes()))
            conn.commit()
        
        # Connect to related topics
        for related_topic in related_topics:
            if related_topic not in self.knowledge_graph:
                if UNIVERSAL_ENCODER:
                    rel_embedding = UNIVERSAL_ENCODER.encode([related_topic])[0]
                else:
                    rel_embedding = np.random.rand(768)  # Fallback
                self.knowledge_graph.add_node(related_topic, embedding=rel_embedding)
                
                # Save to database
                c.execute("INSERT OR REPLACE INTO concepts (concept, embedding) VALUES (?, ?)", 
                         (related_topic, rel_embedding.tobytes()))
            
            if not self.knowledge_graph.has_edge(topic, related_topic):
                # Calculate relation strength based on embedding similarity
                try:
                    sim = np.dot(self.knowledge_graph.nodes[topic]['embedding'], 
                                 self.knowledge_graph.nodes[related_topic]['embedding'])
                except:
                    sim = 0.5
                    
                self.knowledge_graph.add_edge(topic, related_topic, relation="related", strength=sim)
                
                # Save to database
                c.execute("""INSERT OR REPLACE INTO knowledge_graph 
                            (source, target, relation, strength) VALUES (?, ?, ?, ?)""", 
                         (topic, related_topic, "related", float(sim)))
        
        conn.commit()
        conn.close()
    
    def find_related_topics(self, topic, n=5):
        """Find related topics from knowledge graph"""
        if topic not in self.knowledge_graph:
            return []
        
        # Get neighbors with highest strength
        neighbors = list(self.knowledge_graph.neighbors(topic))
        strengths = [self.knowledge_graph[topic][nbr]['strength'] for nbr in neighbors]
        
        # Sort by strength descending
        sorted_indices = np.argsort(strengths)[::-1]
        return [neighbors[i] for i in sorted_indices[:min(n, len(neighbors))]]
    
    def universal_question_generation(self, topic):
        """Generate questions using universal knowledge"""
        # Get related concepts
        related_concepts = self.find_related_topics(topic, 3)
        
        # Generate interdisciplinary questions
        questions = [
            f"What is the fundamental nature of {topic}?",
            f"How does {topic} relate to {related_concepts[0] if related_concepts else 'other fields'}?",
            f"What universal principles underlie {topic}?",
            f"How has {topic} evolved across different domains?",
            f"What are the cross-cutting applications of {topic}?",
            f"What ethical considerations apply universally to {topic}?",
            f"How might {topic} transform in the next decade across all domains?"
        ]
        
        return questions
    
    def voice_input(self):
        """Capture voice input"""
        try:
            with sr.Microphone() as source:
                audio = self.recognizer.listen(source, timeout=5)
                text = self.recognizer.recognize_google(audio)
                return text
        except sr.UnknownValueError:
            raise Exception("Could not understand audio")
        except sr.RequestError as e:
            raise Exception(f"Speech recognition error: {str(e)}")
        except Exception as e:
            raise Exception(f"Voice input failed: {str(e)}")
    
    def text_to_speech(self, text):
        """Convert text to speech"""
        try:
            self.engine.say(text)
            self.engine.runAndWait()
        except Exception as e:
            logger.error(f"Speech synthesis failed: {str(e)}")
    
    def create_presentation(self, results):
        """Generate PowerPoint presentation from results"""
        from pptx import Presentation
        from pptx.util import Inches
        
        try:
            prs = Presentation()
            
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"Research: {results['topic']}"
            subtitle.text = f"Generated by HJKAI on {datetime.now().strftime('%Y-%m-%d')}"
            
            # Summary slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Executive Summary"
            content.text = results.get("summary", "No summary available")
            
            # Questions slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Research Questions"
            content.text = "\n".join(results.get("questions", []))
            
            # Save presentation
            filepath = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx")]
            )
            if filepath:
                prs.save(filepath)
                return filepath
        except ImportError:
            logger.error("python-pptx not installed")
            return None
        except Exception as e:
            logger.error(f"Presentation creation failed: {str(e)}")
            return None

# --- GUI with Universal Mode ---
class HJKAIGUI:
    def __init__(self):
        self.core = HJKAICore()
        self.search_engine = "google"
        self.theme_mode = "dark"
        self.universal_mode = False
        self.setup_gui()
        
    def setup_gui(self):
        ctk.set_appearance_mode("dark")
        ctk.set_default_color_theme("blue")

        self.app = ctk.CTk()
        self.app.geometry("1200x900")
        self.app.title("HJKAI - Advanced AI Research Assistant")
        self.app.grid_columnconfigure(0, weight=1)
        self.app.grid_rowconfigure(0, weight=1)

        # Create main frame
        main_frame = ctk.CTkFrame(self.app)
        main_frame.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")
        main_frame.grid_columnconfigure(0, weight=1)
        main_frame.grid_rowconfigure(1, weight=1)

        # Create tabs
        self.tabs = ctk.CTkTabview(main_frame)
        self.tabs.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")
        self.tabs.grid_columnconfigure(0, weight=1)
        self.tabs.grid_rowconfigure(0, weight=1)
        
        # Add tabs
        self.main_tab = self.tabs.add("AI Processor")
        self.reason_tab = self.tabs.add("Reasoning System")
        self.kg_tab = self.tabs.add("Knowledge Graph")
        self.history_tab = self.tabs.add("History")
        
        # Setup tabs
        self.setup_main_tab()
        self.setup_reason_tab()
        self.setup_kg_tab()
        self.setup_history_tab()
        
        # Status bar
        self.status = ctk.CTkLabel(self.app, text="Ready", anchor="w")
        self.status.grid(row=1, column=0, sticky="ew", padx=10, pady=5)
        
        # Bind close event
        self.app.protocol("WM_DELETE_WINDOW", self.on_close)
        self.app.mainloop()
    
    def setup_main_tab(self):
        # Input frame
        input_frame = ctk.CTkFrame(self.main_tab)
        input_frame.pack(pady=10, padx=10, fill="x")
        
        # Topic input
        self.entry = ctk.CTkEntry(input_frame, placeholder_text="Enter research topic...")
        self.entry.pack(side="left", padx=5, fill="x", expand=True)
        
        # Voice input button
        voice_btn = ctk.CTkButton(input_frame, text="🎤", width=40, command=self.voice_input)
        voice_btn.pack(side="left", padx=5)
        
        # Universal mode toggle
        self.universal_toggle = ctk.CTkSwitch(input_frame, text="Universal Mode", command=self.toggle_universal_mode)
        self.universal_toggle.pack(side="left", padx=10)
        
        # Run button
        self.run_button = ctk.CTkButton(input_frame, text="Run Analysis", command=self.start_ai_thread)
        self.run_button.pack(side="right", padx=5)
        
        # Output tabs
        output_tabs = ctk.CTkTabview(self.main_tab)
        output_tabs.pack(pady=10, padx=10, fill="both", expand=True)
        
        self.results_tab = output_tabs.add("Results")
        self.questions_tab = output_tabs.add("Questions")
        self.summary_tab = output_tabs.add("Summary")
        self.sentiment_tab = output_tabs.add("Sentiment")
        
        # Results tab
        self.output_box = ctk.CTkTextbox(self.results_tab, font=("Consolas", 12))
        self.output_box.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Questions tab
        self.questions_box = ctk.CTkTextbox(self.questions_tab, font=("Arial", 14))
        self.questions_box.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Summary tab
        self.summary_box = ctk.CTkTextbox(self.summary_tab, font=("Arial", 14), wrap="word")
        self.summary_box.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Sentiment tab
        sentiment_frame = ctk.CTkFrame(self.sentiment_tab)
        sentiment_frame.pack(pady=10, padx=10, fill="both", expand=True)
        
        self.sentiment_chart_frame = ctk.CTkFrame(sentiment_frame, height=300)
        self.sentiment_chart_frame.pack(pady=10, padx=10, fill="both", expand=True)
        
        self.sentiment_table = ctk.CTkTextbox(sentiment_frame, height=200)
        self.sentiment_table.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Progress bar
        self.progress = ctk.CTkProgressBar(self.main_tab, mode="indeterminate")
        self.progress.pack(pady=5, padx=10, fill="x")
        self.progress.set(0)
        
        # Action buttons
        btn_frame = ctk.CTkFrame(self.main_tab)
        btn_frame.pack(pady=10, padx=10, fill="x")
        
        save_btn = ctk.CTkButton(btn_frame, text="Save Report", command=self.save_results)
        save_btn.pack(side="left", padx=5)
        
        ppt_btn = ctk.CTkButton(btn_frame, text="Create PPT", command=self.create_ppt)
        ppt_btn.pack(side="left", padx=5)
        
        speak_btn = ctk.CTkButton(btn_frame, text="Speak Summary", command=self.speak_summary)
        speak_btn.pack(side="right", padx=5)
    
    def setup_reason_tab(self):
        reason_frame = ctk.CTkFrame(self.reason_tab)
        reason_frame.pack(pady=10, padx=10, fill="both", expand=True)
        
        self.reason_box = ctk.CTkTextbox(reason_frame, font=("Consolas", 12))
        self.reason_box.pack(pady=10, padx=10, fill="both", expand=True)
        self.reason_box.insert("1.0", self.core.reasoner.system_report())
        
        # Refresh button
        refresh_btn = ctk.CTkButton(reason_frame, text="Refresh", command=self.update_reasoning)
        refresh_btn.pack(pady=5)
    
    def setup_kg_tab(self):
        kg_frame = ctk.CTkFrame(self.kg_tab)
        kg_frame.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Visualization frame
        self.kg_frame = ctk.CTkFrame(kg_frame, height=500)
        self.kg_frame.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Controls
        btn_frame = ctk.CTkFrame(kg_frame)
        btn_frame.pack(pady=10, fill="x")
        
        refresh_btn = ctk.CTkButton(btn_frame, text="Refresh Graph", command=self.visualize_knowledge_graph)
        refresh_btn.pack(side="left", padx=5)
        
        export_btn = ctk.CTkButton(btn_frame, text="Export", command=self.export_knowledge_graph)
        export_btn.pack(side="right", padx=5)
        
        # Concept explorer
        concept_frame = ctk.CTkFrame(kg_frame)
        concept_frame.pack(pady=10, fill="both", expand=True)
        
        ctk.CTkLabel(concept_frame, text="Concept Explorer:").pack(anchor="w")
        
        search_frame = ctk.CTkFrame(concept_frame)
        search_frame.pack(fill="x", pady=5)
        
        self.concept_entry = ctk.CTkEntry(search_frame, placeholder_text="Enter concept")
        self.concept_entry.pack(side="left", padx=5, fill="x", expand=True)
        
        search_btn = ctk.CTkButton(search_frame, text="Explore", command=self.explore_concept)
        search_btn.pack(side="left", padx=5)
        
        self.concept_details = ctk.CTkTextbox(concept_frame, wrap="word")
        self.concept_details.pack(fill="both", expand=True, pady=5)
    
    def setup_history_tab(self):
        history_frame = ctk.CTkFrame(self.history_tab)
        history_frame.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Treeview for history
        columns = ("id", "topic", "timestamp")
        self.history_tree = ttk.Treeview(history_frame, columns=columns, show="headings")
        
        # Define headings
        self.history_tree.heading("id", text="ID")
        self.history_tree.heading("topic", text="Topic")
        self.history_tree.heading("timestamp", text="Timestamp")
        
        # Set column widths
        self.history_tree.column("id", width=50)
        self.history_tree.column("topic", width=300)
        self.history_tree.column("timestamp", width=150)
        
        self.history_tree.pack(pady=10, padx=10, fill="both", expand=True)
        
        # Load history button
        load_btn = ctk.CTkButton(history_frame, text="Load Selected", command=self.load_history)
        load_btn.pack(pady=5)
        
        # Delete button
        del_btn = ctk.CTkButton(history_frame, text="Delete", command=self.delete_history)
        del_btn.pack(pady=5)
        
        # Load initial history
        self.load_history_data()
    
    def toggle_universal_mode(self):
        self.universal_mode = self.universal_toggle.get()
        mode = "Universal" if self.universal_mode else "Standard"
        self.update_status(f"Switched to {mode} Mode")
    
    def voice_input(self):
        try:
            text = self.core.voice_input()
            self.entry.delete(0, "end")
            self.entry.insert(0, text)
            self.update_status("Voice input captured")
        except Exception as e:
            self.show_error(str(e))
    
    def start_ai_thread(self):
        topic = self.entry.get().strip()
        if not topic:
            self.show_error("Please enter a topic")
            return
            
        # Disable buttons during processing
        self.run_button.configure(state="disabled")
        self.update_status("Processing... (This may take 10-20 seconds)")
        self.output_box.delete("1.0", "end")
        self.output_box.insert("1.0", f"🤖 Processing '{topic}' with HJKAI...\n")
        self.progress.start()
        
        # Start background thread
        threading.Thread(target=self.run_ai, args=(topic,), daemon=True).start()
    
    def run_ai(self, topic):
        try:
            # Update progress
            self.update_progress(0.1, "Starting analysis...")
            
            # Fetch data
            self.update_progress(0.2, f"Fetching data from {self.search_engine}...")
            try:
                raw_data = self.core.fetch_web_data(topic, self.search_engine)
            except Exception as e:
                self.show_error(f"Scraping failed: {str(e)}")
                return
                
            if not raw_data:
                self.show_error("No data found. Try different search terms")
                return
                
            # Update banned words
            self.update_progress(0.3, "Updating knowledge base...")
            clean = [d for d in raw_data if any(word in d.lower() for word in ['science', 'data', 'technology'])]
            spam = [d for d in raw_data if any(word in d.lower() for word in ['buy', 'offer', 'cheap', 'free'])]
            self.core.upgrade_banned_words(clean, spam)
            
            # Filter data
            self.update_progress(0.4, "Filtering content...")
            filtered_data = []
            blocked_data = []
            for item in raw_data:
                passed, reasons, confidence, risk = self.core.reasoner.explain_decision(item)
                if passed:
                    filtered_data.append(item)
                else:
                    blocked_data.append({
                        "text": item, 
                        "reasons": reasons,
                        "confidence": confidence,
                        "risk_level": risk
                    })
            
            # Sentiment analysis
            self.update_progress(0.5, "Analyzing sentiment...")
            sentiments, sentiment_counts = self.core.analyze_sentiment(filtered_data)
            
            # Train and cluster
            self.update_progress(0.6, "Clustering content...")
            clustered, keywords, labels, topics, summary = self.core.train_on_data(filtered_data)
            
            # Build knowledge graph
            self.update_progress(0.7, "Building knowledge graph...")
            knowledge_graph = self.core.build_knowledge_graph(filtered_data, keywords)
            
            # Universal knowledge integration
            related_topics = []
            if self.universal_mode:
                self.update_progress(0.75, "Connecting universal knowledge...")
                # Find related topics
                related_topics = self.core.find_related_topics(topic, 5)
                # Update knowledge graph
                self.core.update_knowledge_graph(topic, related_topics)
            
            # Generate questions
            self.update_progress(0.8, "Generating research questions...")
            if self.universal_mode:
                questions = self.core.universal_question_generation(topic)
            else:
                try:
                    questions = self.core.ml_generate_questions(topic)
                except:
                    questions = self.core.generate_default_questions(topic)
            
            # Auto-improve system
            self.update_progress(0.9, "Optimizing system...")
            self.core.auto_improve(filtered_data)
            
            # Prepare results
            results = {
                "topic": topic,
                "questions": questions,
                "accepted_data": filtered_data,
                "blocked_data": blocked_data,
                "clusters": clustered,
                "keywords": keywords,
                "banned_words": banned_words,
                "sentiments": sentiments,
                "sentiment_distribution": sentiment_counts,
                "knowledge_graph": list(knowledge_graph.edges(data=True)),
                "topic_modeling": topics,
                "summary": summary,
                "universal_mode": self.universal_mode,
                "related_topics": related_topics,
                "timestamp": datetime.now().isoformat()
            }
            
            # Save to history
            self.save_to_history(topic, results)
            
            # Update GUI
            self.app.after(0, lambda: self.update_ui(results))
            self.update_progress(1.0, "Processing complete!")
            
        except Exception as e:
            self.show_error(f"Processing error: {str(e)}")
            logger.exception("Processing failed")
        finally:
            self.enable_buttons()
            self.progress.stop()
    
    def update_ui(self, results):
        """Update all UI elements with results"""
        # Format results for display
        formatted_results = json.dumps(results, indent=4)
        
        # Update output boxes
        self.output_box.delete("1.0", "end")
        self.output_box.insert("1.0", formatted_results)
        
        self.questions_box.delete("1.0", "end")
        self.questions_box.insert("1.0", "\n".join(results["questions"]))
        
        self.summary_box.delete("1.0", "end")
        self.summary_box.insert("1.0", results["summary"])
        
        # Update sentiment display
        self.show_sentiment(results["sentiment_distribution"], results["sentiments"])
        
        # Update knowledge graph visualization
        self.visualize_knowledge_graph()
        
        # Update status
        if results["universal_mode"]:
            related = ", ".join(results["related_topics"]) if results["related_topics"] else "None found"
            self.update_status(f"✅ Universal analysis complete! Related topics: {related}")
        else:
            self.update_status("✅ Analysis complete!")
    
    def show_sentiment(self, distribution, sentiments):
        """Update sentiment visualization"""
        # Clear previous content
        for widget in self.sentiment_chart_frame.winfo_children():
            widget.destroy()
        
        # Create pie chart
        fig, ax = plt.subplots(figsize=(5, 3))
        labels = list(distribution.keys())
        sizes = list(distribution.values())
        ax.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90)
        ax.axis('equal')
        ax.set_title("Sentiment Distribution")
        
        # Embed in GUI
        canvas = FigureCanvasTkAgg(fig, master=self.sentiment_chart_frame)
        canvas.draw()
        canvas.get_tk_widget().pack(fill="both", expand=True)
        
        # Create sentiment table
        table_text = "Text\t\tSentiment\tScore\n"
        table_text += "-"*80 + "\n"
        for sent in sentiments[:10]:  # Show first 10
            truncated = sent['text'][:100] + "..." if len(sent['text']) > 100 else sent['text']
            table_text += f"{truncated}\t{sent['sentiment']}\t{sent['score']:.2f}\n"
        
        self.sentiment_table.delete("1.0", "end")
        self.sentiment_table.insert("1.0", table_text)
    
    def visualize_knowledge_graph(self):
        """Visualize the knowledge graph"""
        # Clear previous content
        for widget in self.kg_frame.winfo_children():
            widget.destroy()
        
        if not self.core.knowledge_graph.nodes:
            label = ctk.CTkLabel(self.kg_frame, text="Knowledge graph is empty")
            label.pack(pady=50)
            return
            
        try:
            # Create plot
            fig, ax = plt.subplots(figsize=(10, 8))
            
            # Draw graph
            pos = nx.spring_layout(self.core.knowledge_graph, seed=42)
            nx.draw_networkx_nodes(self.core.knowledge_graph, pos, node_size=500, ax=ax)
            nx.draw_networkx_edges(self.core.knowledge_graph, pos, width=1.0, alpha=0.5, ax=ax)
            nx.draw_networkx_labels(self.core.knowledge_graph, pos, font_size=10, font_family="sans-serif", ax=ax)
            
            ax.set_title("Knowledge Graph")
            
            # Embed in GUI
            canvas = FigureCanvasTkAgg(fig, master=self.kg_frame)
            canvas.draw()
            canvas.get_tk_widget().pack(fill="both", expand=True)
        except Exception as e:
            logger.error(f"Graph visualization failed: {str(e)}")
            label = ctk.CTkLabel(self.kg_frame, text=f"Graph error: {str(e)}")
            label.pack(pady=50)
    
    def explore_concept(self):
        concept = self.concept_entry.get().strip()
        if not concept:
            return
            
        details = f"Concept: {concept}\n\n"
        
        if concept in self.core.knowledge_graph:
            node = self.core.knowledge_graph.nodes[concept]
            details += f"Connections: {self.core.knowledge_graph.degree(concept)}\n"
            
            # Show connections
            neighbors = list(self.core.knowledge_graph.neighbors(concept))
            strengths = [self.core.knowledge_graph[concept][nbr]['strength'] for nbr in neighbors]
            
            details += "\nStrongest Connections:\n"
            # Get top 5 connections
            sorted_indices = np.argsort(strengths)[::-1][:5]
            for i in sorted_indices:
                neighbor = neighbors[i]
                strength = strengths[i]
                details += f"- {neighbor} (strength: {strength:.2f})\n"
        else:
            details += "Concept not in knowledge base\n"
            details += "You can add it by analyzing a related topic"
        
        self.concept_details.delete("1.0", "end")
        self.concept_details.insert("1.0", details)
    
    def save_to_history(self, topic, results):
        """Save results to history database"""
        # Save results to file
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{topic.replace(' ', '_')}_{timestamp}.json"
        with open(filename, "w") as f:
            json.dump(results, f)
        
        # Save to database
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        c.execute("INSERT INTO queries (topic, results_file) VALUES (?, ?)", 
                 (topic, filename))
        conn.commit()
        conn.close()
        
        # Reload history
        self.load_history_data()
    
    def load_history_data(self):
        """Load history data into treeview"""
        # Clear existing data
        for item in self.history_tree.get_children():
            self.history_tree.delete(item)
        
        # Fetch history
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        c.execute("SELECT id, topic, timestamp FROM queries ORDER BY timestamp DESC")
        for row in c.fetchall():
            self.history_tree.insert("", "end", values=row)
        conn.close()
    
    def load_history(self):
        """Load selected history item"""
        selection = self.history_tree.selection()
        if not selection:
            return
            
        item = self.history_tree.item(selection[0])
        filepath = item['values'][2]  # Results_file is the third column
        
        try:
            with open(filepath, "r") as f:
                results = json.load(f)
                self.update_ui(results)
                self.update_status(f"Loaded history: {item['values'][1]}")
        except Exception as e:
            self.show_error(f"Failed to load history: {str(e)}")
    
    def delete_history(self):
        """Delete selected history item"""
        selection = self.history_tree.selection()
        if not selection:
            return
            
        item = self.history_tree.item(selection[0])
        query_id = item['values'][0]
        
        # Delete from database and filesystem
        conn = sqlite3.connect('hjkai.db')
        c = conn.cursor()
        c.execute("SELECT results_file FROM queries WHERE id = ?", (query_id,))
        result = c.fetchone()
        if result:
            try:
                os.remove(result[0])
            except:
                pass
        c.execute("DELETE FROM queries WHERE id = ?", (query_id,))
        conn.commit()
        conn.close()
        
        # Reload history
        self.load_history_data()
        self.update_status(f"Deleted history item {query_id}")
    
    def save_results(self):
        """Save current results to file"""
        data = self.output_box.get("1.0", "end-1c")
        filepath = filedialog.asksaveasfilename(
            defaultextension=".json",
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")]
        )
        if filepath:
            with open(filepath, "w") as f:
                f.write(data)
            self.update_status(f"💾 Results saved to {filepath}")
    
    def create_ppt(self):
        """Create PowerPoint from results"""
        data = self.output_box.get("1.0", "end-1c")
        try:
            results = json.loads(data)
        except:
            self.show_error("Invalid results format")
            return
            
        filepath = self.core.create_presentation(results)
        if filepath:
            self.update_status(f"💾 Presentation saved to {filepath}")
        else:
            self.show_error("Failed to create presentation")
    
    def speak_summary(self):
        """Speak the summary aloud"""
        summary = self.summary_box.get("1.0", "end-1c")
        if summary:
            threading.Thread(target=self.core.text_to_speech, args=(summary,), daemon=True).start()
            self.update_status("Speaking summary...")
    
    def export_knowledge_graph(self):
        """Export knowledge graph to file"""
        filepath = filedialog.asksaveasfilename(
            defaultextension=".graphml",
            filetypes=[("GraphML files", "*.graphml"), ("All files", "*.*")]
        )
        if filepath:
            try:
                nx.write_graphml(self.core.knowledge_graph, filepath)
                self.update_status(f"💾 Knowledge graph exported to {filepath}")
            except Exception as e:
                self.show_error(f"Export failed: {str(e)}")
    
    def update_reasoning(self):
        """Update the reasoning display"""
        self.reason_box.delete("1.0", "end")
        self.reason_box.insert("1.0", self.core.reasoner.system_report())
        self.update_status("Reasoning system refreshed")
    
    def update_progress(self, value, message):
        """Thread-safe progress update"""
        self.app.after(0, lambda: self.progress.set(value))
        self.app.after(0, lambda: self.status.configure(text=message))
    
    def update_status(self, message):
        """Thread-safe status update"""
        self.app.after(0, lambda: self.status.configure(text=message))
    
    def show_error(self, message):
        """Thread-safe error display"""
        self.app.after(0, lambda: messagebox.showerror("Error", message))
        self.app.after(0, lambda: self.status.configure(text=f"❌ {message}"))
    
    def enable_buttons(self):
        """Re-enable UI buttons"""
        self.app.after(0, lambda: self.run_button.configure(state="normal"))
    
    def on_close(self):
        if messagebox.askokcancel("Quit", "Do you want to quit?"):
            self.app.destroy()
            sys.exit()

if __name__ == "__main__":
    HJKAIGUI()
